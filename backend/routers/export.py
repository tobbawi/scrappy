import io
import json
from datetime import datetime
from typing import Optional

from fastapi import APIRouter, Depends, Query, Response
from sqlmodel import Session, select

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

from database import get_session
from models import ReferenceCase
from routers.cases import _case_filters

router = APIRouter(prefix="/api/export", tags=["export"])

# Colour palette
_HEADER_BG = RGBColor(0x4F, 0x46, 0xE5)
_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
_LABEL_FG  = RGBColor(0x6B, 0x72, 0x80)
_BODY_FG   = RGBColor(0x11, 0x18, 0x27)
_QUOTE_BG  = RGBColor(0xEE, 0xF2, 0xFF)
_QUOTE_FG  = RGBColor(0x31, 0x2E, 0x81)
_FOOTER_BG = RGBColor(0xF9, 0xFA, 0xFB)
_FOOTER_FG = RGBColor(0x6B, 0x72, 0x80)

# Slide dimensions (EMU, 16:9 widescreen)
SLIDE_W = 12_192_000
SLIDE_H =  6_858_000

# Row heights (EMU)
H_HEADER  = 1_143_000
H_CH_SOL  = 2_286_000
H_RESULTS = 1_371_600
H_QUOTE   = 1_143_000
H_FOOTER  =   914_400

# Padding (EMU, ~0.15 in)
PAD = 137_160


def _trunc(s: Optional[str], n: int) -> str:
    if not s:
        return ""
    return s[:n] + ("…" if len(s) > n else "")


def _parse_list(raw: Optional[str]) -> list:
    if not raw:
        return []
    try:
        val = json.loads(raw)
        if isinstance(val, list):
            return [str(x) for x in val]
    except (json.JSONDecodeError, TypeError):
        pass
    return [x.strip() for x in raw.split(",") if x.strip()]


def _add_bg(slide, left: int, top: int, width: int, height: int, color: RGBColor) -> None:
    """Add a filled rectangle with no border as a background band."""
    shape = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _txbox(slide, left: int, top: int, width: int, height: int):
    """Add a textbox and return its text_frame."""
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _run(para, text: str, size: float, bold: bool = False, italic: bool = False,
         color: RGBColor = _BODY_FG):
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return run


def _label_body(slide, left: int, top: int, width: int, height: int,
                label: str, body: str) -> None:
    """Render a section with a small caps label and a body paragraph."""
    tf = _txbox(slide, left, top, width, height)
    _run(tf.paragraphs[0], label, size=8, color=_LABEL_FG)
    p2 = tf.add_paragraph()
    _run(p2, body, size=10, color=_BODY_FG)


def _build_slide(prs: Presentation, case: ReferenceCase) -> None:
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    top = 0

    # ── Header ──────────────────────────────────────────────────────────────
    _add_bg(slide, 0, top, SLIDE_W, H_HEADER, _HEADER_BG)

    half_w = SLIDE_W // 2

    # Left: company ID
    tf_co = _txbox(slide, PAD, top + PAD, half_w - PAD * 2, H_HEADER // 2)
    _run(tf_co.paragraphs[0], case.company_id or "", size=10, color=_WHITE)

    # Left: metadata row
    meta_parts = [p for p in [
        case.customer_industry,
        case.customer_country,
        str(case.publish_date)[:10] if case.publish_date else None,
    ] if p]
    tf_meta = _txbox(slide, PAD, top + H_HEADER // 2, half_w - PAD * 2, H_HEADER // 2 - PAD)
    _run(tf_meta.paragraphs[0], "  •  ".join(meta_parts), size=9, color=_WHITE)

    # Right: customer name (large, bold)
    tf_cust = _txbox(slide, half_w, top + PAD, half_w - PAD, H_HEADER - PAD * 2)
    _run(tf_cust.paragraphs[0], case.customer_name or "", size=20, bold=True, color=_WHITE)

    top += H_HEADER

    # ── Challenge + Solution ─────────────────────────────────────────────────
    has_ch  = bool(case.challenge)
    has_sol = bool(case.solution)
    if has_ch or has_sol:
        col_w = SLIDE_W // 2
        _label_body(slide, PAD, top + PAD, col_w - PAD * 2, H_CH_SOL - PAD * 2,
                    "CHALLENGE", _trunc(case.challenge, 300))
        _label_body(slide, col_w + PAD, top + PAD, col_w - PAD * 2, H_CH_SOL - PAD * 2,
                    "SOLUTION", _trunc(case.solution, 300))
        top += H_CH_SOL

    # ── Results ──────────────────────────────────────────────────────────────
    if case.results:
        _label_body(slide, PAD, top + PAD, SLIDE_W - PAD * 2, H_RESULTS - PAD * 2,
                    "RESULTS", _trunc(case.results, 400))
        top += H_RESULTS

    # ── Quote ────────────────────────────────────────────────────────────────
    if case.quote:
        _add_bg(slide, 0, top, SLIDE_W, H_QUOTE, _QUOTE_BG)
        attribution = ""
        if case.quote_author:
            attribution = f"  \u2014 {case.quote_author}"
            if case.quote_author_company:
                attribution += f", {case.quote_author_company}"
        quote_text = f"\u201c{_trunc(case.quote, 200)}\u201d{attribution}"
        tf_q = _txbox(slide, PAD * 2, top + PAD, SLIDE_W - PAD * 4, H_QUOTE - PAD * 2)
        _run(tf_q.paragraphs[0], quote_text, size=11, italic=True, color=_QUOTE_FG)
        top += H_QUOTE

    # ── Footer ───────────────────────────────────────────────────────────────
    _add_bg(slide, 0, top, SLIDE_W, H_FOOTER, _FOOTER_BG)
    tags     = ", ".join(_parse_list(case.tags)[:6]) or "\u2014"
    products = ", ".join(_parse_list(case.products_used)[:6]) or "\u2014"
    footer_text = f"Tags: {tags}  \u2502  Products: {products}  \u2502  {case.url or ''}"
    tf_f = _txbox(slide, PAD, top + PAD, SLIDE_W - PAD * 2, H_FOOTER - PAD * 2)
    _run(tf_f.paragraphs[0], footer_text, size=8, color=_FOOTER_FG)


def build_pptx(cases: list) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    for case in cases:
        _build_slide(prs, case)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


@router.get("/pptx")
def export_pptx(
    company:  Optional[str] = None,
    industry: Optional[str] = None,
    country:  Optional[str] = None,
    q:        Optional[str] = Query(None, max_length=200),
    since:    Optional[str] = None,
    new_only: bool = False,
    sort:     str  = Query("first_seen", pattern="^(first_seen|publish_date)$"),
    session:  Session = Depends(get_session),
):
    since_dt: Optional[datetime] = None
    if since:
        try:
            since_dt = datetime.fromisoformat(since)
        except ValueError:
            pass

    clauses  = _case_filters(company, industry, country, q, since_dt, new_only)
    sort_col = ReferenceCase.first_seen if sort == "first_seen" else ReferenceCase.publish_date
    cases    = session.exec(
        select(ReferenceCase).where(*clauses).order_by(sort_col.desc())
    ).all()

    buf = build_pptx(cases)
    return Response(
        content=buf.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": 'attachment; filename="scrappy-cases.pptx"'},
    )
